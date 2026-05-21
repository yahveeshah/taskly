from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Image folder path
IMAGE_FOLDER = r"C:\Users\Admin\Pictures\Screenshots\taskly_images"
OUTPUT_PATH = r"C:\Users\Admin\Herd\project\Taskly_Presentation.pptx"

# Colors
NAVY = RgbColor(0, 0, 128)
WHITE = RgbColor(255, 255, 255)
DARK_GREY = RgbColor(45, 45, 45)

# Slide dimensions
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

def add_header_bar(slide, title):
    """Add navy header bar with white title text to slide"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.color = NAVY
    header.line.fill.background()
    
    text_frame = header.text_frame
    text_frame.text = title
    text_frame.paragraphs[0].font.size = Pt(28)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = WHITE
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)

def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add a text box to the slide"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.alignment = alignment
    return text_box

def create_presentation():
    # Create presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Get all image files
    all_images = [f for f in os.listdir(IMAGE_FOLDER) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]
    
    # Separate logos and screenshots
    logos = [f for f in all_images if any(keyword in f.lower() for keyword in ['final', 'logo', 'images'])]
    screenshots = [f for f in all_images if f not in logos]
    screenshots.sort()  # Sort screenshots alphabetically
    
    # SLIDE 1 - Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.color = NAVY
    
    # Add logos side by side at top center
    if len(logos) >= 2:
        logo1_path = os.path.join(IMAGE_FOLDER, logos[0])
        logo2_path = os.path.join(IMAGE_FOLDER, logos[1])
        slide1.shapes.add_picture(logo1_path, Inches(5.5), Inches(0.3), width=Inches(1.5))
        slide1.shapes.add_picture(logo2_path, Inches(7.2), Inches(0.3), width=Inches(1.5))
    elif len(logos) == 1:
        logo_path = os.path.join(IMAGE_FOLDER, logos[0])
        slide1.shapes.add_picture(logo_path, Inches(5.9), Inches(0.3), width=Inches(1.5))
    
    # Title
    add_text_box(slide1, "Taskly — Task Management Web Application", 
                 Inches(1), Inches(2.2), Inches(11.33), Inches(0.8), 
                 font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Student info
    info_text = (
        "Student Name: Yahvee Shah\n"
        "Enrollment Number: 24012250210139\n"
        "Branch: Computer Engineering\n"
        "College Name: LJ Polytechnic\n"
        "Internship Duration: 4 weeks\n"
        "Technology Stack: PHP 8.4, Laravel 11, SQLite, Blade, SortableJS, Vanilla JS, CSS3, Git and GitHub"
    )
    add_text_box(slide1, info_text, Inches(2), Inches(3.2), Inches(9.33), Inches(3), 
                 font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # SLIDE 2 - Problem Statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.color = WHITE
    add_header_bar(slide2, "Problem Statement")
    
    problem_text = (
        "Students and professionals struggle to manage tasks efficiently due to lack of structured, "
        "visually appealing and role based productivity tools. Taskly solves this by providing "
        "priority based organization, deadline tracking, progress visualization, team collaboration "
        "with role based access and a personalized dashboard within a dual theme interface."
    )
    add_text_box(slide2, problem_text, Inches(1), Inches(1.2), Inches(11.33), Inches(5.5), 
                 font_size=20, color=NAVY)
    
    # SLIDE 3 - Technology Stack
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.color = WHITE
    add_header_bar(slide3, "Technology Stack")
    
    tech_items = [
        "PHP 8.4 — server side scripting language",
        "Laravel 11 — PHP web application framework",
        "SQLite — lightweight relational database",
        "Blade Templating — Laravel view engine",
        "SortableJS — drag and drop library",
        "Vanilla JavaScript — client side interactivity",
        "CSS3 Custom Properties — dual theme design system",
        "Git and GitHub — version control and deployment"
    ]
    
    # Left column
    left_text = "\n".join(tech_items[:4])
    add_text_box(slide3, left_text, Inches(1), Inches(1.2), Inches(5.5), Inches(5.5), 
                 font_size=18, color=NAVY)
    
    # Right column
    right_text = "\n".join(tech_items[4:])
    add_text_box(slide3, right_text, Inches(6.83), Inches(1.2), Inches(5.5), Inches(5.5), 
                 font_size=18, color=NAVY)
    
    # SLIDE 4 - Features and Functionalities
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.color = WHITE
    add_header_bar(slide4, "Features and Functionalities")
    
    features = [
        "User Authentication and Registration",
        "Role Based Access Control",
        "Personal and Team Mode",
        "Task CRUD with Priority and Deadline",
        "Drag and Drop Task Board",
        "Progress Tracking Page",
        "Bar Graph Visualization",
        "Deadline Calendar Widget",
        "Streak Tracker",
        "Team Management with Team Codes",
        "Manager Member Drilldown",
        "Dual Theme System Lavender and Vintage",
        "Collapsible Sidebar",
        "Task Timeline Table",
        "Dashboard with Weekly Summary and Greeting"
    ]
    
    # Left column
    left_features = "\n".join(features[:8])
    add_text_box(slide4, left_features, Inches(1), Inches(1.2), Inches(5.5), Inches(5.5), 
                 font_size=16, color=NAVY)
    
    # Right column
    right_features = "\n".join(features[8:])
    add_text_box(slide4, right_features, Inches(6.83), Inches(1.2), Inches(5.5), Inches(5.5), 
                 font_size=16, color=NAVY)
    
    # SLIDE 5 - Internship Tasks Completed
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.color = WHITE
    add_header_bar(slide5, "Tasks Completed During Internship")
    
    tasks = [
        "1. Set up Laravel project with Herd",
        "2. Designed database schema and migrations",
        "3. Built user authentication system",
        "4. Implemented role based access control",
        "5. Built task CRUD with priority and deadline",
        "6. Designed three column priority board",
        "7. Implemented drag and drop using SortableJS",
        "8. Built progress tracking page with encouraging cards",
        "9. Built bar graph visualization",
        "10. Built deadline calendar widget",
        "11. Built streak tracker",
        "12. Built team management system with team codes",
        "13. Implemented dual theme system Lavender and Vintage",
        "14. Built collapsible sidebar",
        "15. Added smooth animations and page transitions",
        "16. Built professional task timeline table",
        "17. Deployed project on GitHub"
    ]
    
    tasks_text = "\n".join(tasks)
    add_text_box(slide5, tasks_text, Inches(1), Inches(1.2), Inches(11.33), Inches(5.5), 
                 font_size=16, color=NAVY)
    
    # SLIDE 6 - Skills Gained
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.color = WHITE
    add_header_bar(slide6, "Skills Gained")
    
    technical_skills = (
        "Technical Skills:\n\n"
        "Laravel MVC Architecture\n"
        "Database Design and Migrations\n"
        "Authentication and Authorization with Policies\n"
        "RESTful Routing\n"
        "Blade Templating Engine\n"
        "JavaScript DOM Manipulation\n"
        "CSS Custom Properties and Theming\n"
        "Git and GitHub Version Control"
    )
    
    soft_skills = (
        "Soft Skills:\n\n"
        "Problem Solving\n"
        "UI UX Design Thinking\n"
        "Project Planning\n"
        "Time Management"
    )
    
    add_text_box(slide6, technical_skills, Inches(1), Inches(1.2), Inches(5.5), Inches(5.5), 
                 font_size=16, color=NAVY)
    add_text_box(slide6, soft_skills, Inches(6.83), Inches(1.2), Inches(5.5), Inches(5.5), 
                 font_size=16, color=NAVY)
    
    # SLIDE 7 - Project Screenshots Part 1 (first 7 screenshots)
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    slide7.background.fill.solid()
    slide7.background.fill.fore_color.color = WHITE
    add_header_bar(slide7, "Project Screenshots")
    
    screenshots_part1 = screenshots[:7]
    for i, screenshot in enumerate(screenshots_part1):
        row = i // 2
        col = i % 2
        img_path = os.path.join(IMAGE_FOLDER, screenshot)
        left = Inches(1 + col * 5.8)
        top = Inches(1.2 + row * 2.8)
        slide7.shapes.add_picture(img_path, left, top, width=Inches(5.5))
        # Add label below
        label_text = screenshot.replace('.png', '').replace('.jpg', '')
        add_text_box(slide7, label_text, left, top + Inches(2), Inches(5.5), Inches(0.4), 
                     font_size=10, color=NAVY, alignment=PP_ALIGN.CENTER)
    
    # SLIDE 8 - Project Screenshots Part 2 (remaining 7 screenshots)
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.color = WHITE
    add_header_bar(slide8, "Project Screenshots")
    
    screenshots_part2 = screenshots[7:]
    for i, screenshot in enumerate(screenshots_part2):
        row = i // 2
        col = i % 2
        img_path = os.path.join(IMAGE_FOLDER, screenshot)
        left = Inches(1 + col * 5.8)
        top = Inches(1.2 + row * 2.8)
        slide8.shapes.add_picture(img_path, left, top, width=Inches(5.5))
        # Add label below
        label_text = screenshot.replace('.png', '').replace('.jpg', '')
        add_text_box(slide8, label_text, left, top + Inches(2), Inches(5.5), Inches(0.4), 
                     font_size=10, color=NAVY, alignment=PP_ALIGN.CENTER)
    
    # SLIDE 9 - Code Samples
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    slide9.background.fill.solid()
    slide9.background.fill.fore_color.color = WHITE
    add_header_bar(slide9, "Code Samples")
    
    # Code block 1 - Routes
    code1 = """Route::get('/', [DashboardController::class, 'index'])->name('dashboard');
Route::resource('tasks', TaskController::class);
Route::get('/progress', [ProgressController::class, 'index'])->name('progress');
Route::post('/teams/join', [TeamController::class, 'join'])->name('teams.join');"""
    
    # Code block 2 - Blade Template
    code2 = """<div class="task-card" data-priority="{{ $task->priority }}">
    <h3>{{ $task->title }}</h3>
    <p>{{ $task->description }}</p>
    <span class="badge">{{ $task->priority }}</span>
</div>"""
    
    # Code block 3 - Controller
    code3 = """public function store(Request $request)
{
    $validated = $request->validate([
        'title' => 'required|max:255',
        'priority' => 'required|in:high,medium,low',
    ]);
    return Task::create($validated + ['user_id' => auth()->id()]);
}"""
    
    # Add code blocks with dark grey background
    codes = [code1, code2, code3]
    labels = ["Routes — web.php", "Blade Template — tasks/index.blade.php", "Controller — TaskController.php"]
    
    for i, (code, label) in enumerate(zip(codes, labels)):
        # Add label
        add_text_box(slide9, label, Inches(1 + i * 4), Inches(1.2), Inches(4), Inches(0.4), 
                     font_size=11, bold=True, color=NAVY)
        # Add code background
        code_bg = slide9.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1 + i * 4), Inches(1.6), Inches(4), Inches(4.5)
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.color = DARK_GREY
        code_bg.line.fill.background()
        # Add code text
        code_box = slide9.shapes.add_textbox(Inches(1 + i * 4), Inches(1.6), Inches(4), Inches(4.5))
        text_frame = code_box.text_frame
        text_frame.text = code
        text_frame.word_wrap = True
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.name = "Courier New"
            paragraph.font.color.rgb = WHITE
    
    # SLIDE 10 - Live Demo
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    slide10.background.fill.solid()
    slide10.background.fill.fore_color.color = NAVY
    
    add_text_box(slide10, "Live Demo", Inches(1), Inches(2), Inches(11.33), Inches(0.8), 
                 font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide10, "Project URL: project.test", Inches(1), Inches(3), Inches(11.33), Inches(0.6), 
                 font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    demo_points = (
        "• Task creation with priority and deadline\n"
        "• Drag and drop between priority columns\n"
        "• Theme switching between Lavender and Vintage\n"
        "• Progress tracking and bar graph\n"
        "• Team management with role based access"
    )
    add_text_box(slide10, demo_points, Inches(3), Inches(4), Inches(7.33), Inches(2), 
                 font_size=20, color=WHITE)
    
    # SLIDE 11 - Conclusion
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    slide11.background.fill.solid()
    slide11.background.fill.fore_color.color = WHITE
    add_header_bar(slide11, "Conclusion")
    
    conclusion_text = (
        "Taskly is a fully featured task management web application built using Laravel during a 4 week "
        "internship at LJ Polytechnic. The project demonstrates full stack web development skills including "
        "authentication, role based access control, database design, UI UX design principles, and version "
        "control using Git and GitHub. This internship provided hands on experience building a real world "
        "application from scratch."
    )
    add_text_box(slide11, conclusion_text, Inches(1), Inches(1.2), Inches(11.33), Inches(3.5), 
                 font_size=18, color=NAVY)
    
    add_text_box(slide11, "Thank You", Inches(1), Inches(5), Inches(11.33), Inches(1), 
                 font_size=36, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    
    # Save presentation
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved successfully to: {OUTPUT_PATH}")
    print(f"Total slides created: {len(prs.slides)}")
    print(f"Logos used: {len(logos)}")
    print(f"Screenshots used: {len(screenshots)}")

if __name__ == "__main__":
    create_presentation()
